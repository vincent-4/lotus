"""Slides parser."""

from io import BytesIO
from pathlib import Path

from fsspec import AbstractFileSystem
from llama_index.core.readers.base import BaseReader
from llama_index.core.schema import Document
from llama_index.core.utils import infer_torch_device


class PptxReader(BaseReader):
    """Powerpoint parser.

    Extract text and return it page by page.

    Args:
        should_caption_images (bool): Whether to caption images in the slides.
        caption_model (str): The model to use for image captioning.
        device (str | None): The device to use for image captioning. If None, it will be inferred.
        **gen_kwargs: Keyword arguments to pass to the model for image captioning.
    """

    def __init__(
        self,
        should_caption_images: bool = False,
        caption_model: str = "nlpconnect/vit-gpt2-image-captioning",
        device: str | None = None,
        **gen_kwargs,
    ) -> None:
        try:
            from pptx import Presentation  # noqa
        except ImportError:
            raise ImportError(
                "Please install extra dependencies that are required for "
                "the PptxReader: "
                "`pip install python-pptx`"
            )
        self.should_caption_images = should_caption_images
        if self.should_caption_images:
            self._init_caption_images(caption_model, device)
            self.gen_kwargs = gen_kwargs or {"max_length": 16, "num_beams": 4}

    def _init_caption_images(self, caption_model, device):
        try:
            import torch  # noqa
            from PIL import Image  # noqa
            from transformers import (
                AutoTokenizer,
                VisionEncoderDecoderModel,
                ViTFeatureExtractor,
            )
        except ImportError:
            raise ImportError(
                "Please install extra dependencies that are required for "
                "the PptxReader with Image captions: "
                "`pip install torch transformers python-pptx Pillow`"
            )
        device = device or infer_torch_device()
        self.model = VisionEncoderDecoderModel.from_pretrained(caption_model).to(device)
        self.feature_extractor = ViTFeatureExtractor.from_pretrained(caption_model)
        self.tokenizer = AutoTokenizer.from_pretrained(caption_model)

    def caption_image(self, image_bytes: bytes) -> str:
        """Generate text caption of image."""
        from PIL import Image

        i_image: Image.ImageFile.ImageFile | Image.Image = Image.open(BytesIO(image_bytes))
        if i_image.mode != "RGB":
            i_image = i_image.convert(mode="RGB")

        pixel_values = self.feature_extractor(images=[i_image], return_tensors="pt").pixel_values
        pixel_values = pixel_values.to(self.model.device)

        output_ids = self.model.generate(pixel_values, **self.gen_kwargs)

        preds = self.tokenizer.batch_decode(output_ids, skip_special_tokens=True)
        return preds[0].strip()

    def load_data(
        self,
        file: Path,
        extra_info: dict | None = None,
        fs: AbstractFileSystem | None = None,
    ) -> list[Document]:
        """Parse file."""
        from pptx import Presentation

        def get_shape_text(shape):
            text = f"{shape.text}\n" if hasattr(shape, "text") else ""
            if hasattr(shape, "image") and self.should_caption_images:
                text += f"Image: {self.caption_image(shape.image.blob)}\n\n"
            return text

        if fs:
            with fs.open(file) as f:
                presentation = Presentation(f)
        else:
            presentation = Presentation(file)

        docs = []
        for i, slide in enumerate(presentation.slides):
            text = "".join(get_shape_text(shape) for shape in slide.shapes)
            metadata = {"page_label": i + 1, "file_name": file.name}
            if extra_info is not None:
                metadata.update(extra_info)
            docs.append(Document(text=text, metadata=metadata))

        return docs
